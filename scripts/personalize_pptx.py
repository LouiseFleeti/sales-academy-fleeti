#!/usr/bin/env python3
"""
Génère un PPTX personnalisé à partir des slides PNG extraites du PDF Fleeti.
Chaque slide = image PNG en fond, personnalisations superposées sur la couverture.

Usage:
  python3 personalize_pptx.py --type envoyer --client "Acme" --vehicles 85 \
    --vehicle-types "VL,VUL" --pain-points "tco,silos" --output /tmp/out.pptx
"""

import argparse, io, os, sys, urllib.request
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDES_DIR = os.path.join(os.path.dirname(__file__), '..', 'public')

FLEETI_GREEN = RGBColor(0x1A, 0xB0, 0x8E)
FLEETI_DARK  = RGBColor(0x1A, 0x2E, 0x44)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

PAIN_LABELS = {
    'silos':           'outils en silos',
    'flotte':          'flotte hétérogène',
    'tco':             'TCO non maîtrisé',
    'sinistralite':    'sinistralité',
    'contrats':        'contrats LLD/LOA dispersés',
    'fournisseurs':    'fournisseurs non intégrés',
    'electrification': 'électrification sans visibilité',
    'conformite':      'conformité à risque',
}

VTYPE_SHORT = {
    'VL': 'VL', 'VUL': 'VUL', 'PL': 'PL',
    'engins': 'Engins', 'machines': 'Machines'
}


def add_text_box(slide, left_pct, top_pct, width_pct, height_pct,
                 text, font_size, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False):
    """Ajoute une zone de texte positionnée en % de la slide."""
    W = Emu(9144000)   # 10 inches en EMU
    H = Emu(5143500)   # 5.625 inches en EMU

    txBox = slide.shapes.add_textbox(
        Emu(int(W * left_pct)),
        Emu(int(H * top_pct)),
        Emu(int(W * width_pct)),
        Emu(int(H * height_pct)),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_logo(slide, logo_path, left_pct=0.82, top_pct=0.04, width_pct=0.14):
    try:
        logo_w    = Emu(int(9144000 * width_pct))
        logo_left = Emu(int(9144000 * left_pct))
        logo_top  = Emu(int(5143500 * top_pct))
        slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_w)
    except Exception as e:
        print(f'Warning: logo ignoré ({e})', file=sys.stderr)


def build_pptx(slide_type, client_name, vehicles, vehicle_types, pain_points,
               logo_path, sales_name, output_path):

    slides_folder = os.path.join(SLIDES_DIR, f'slides-{slide_type}')
    slide_files = sorted(f for f in os.listdir(slides_folder) if f.endswith('.png'))

    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10"
    prs.slide_height = Emu(5143500)   # 5.625"

    blank_layout = prs.slide_layouts[6]  # Layout vide

    for idx, fname in enumerate(slide_files):
        slide = prs.slides.add_slide(blank_layout)
        img_path = os.path.join(slides_folder, fname)

        # Image plein fond
        slide.shapes.add_picture(
            img_path, 0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        # ── Slide pain points : index 2 pour rdv (slide 3), index 3 pour envoyer (slide 4)
        pain_slide_idx = 3 if slide_type == 'envoyer' else 2
        if idx == pain_slide_idx and pain_points:
            from pptx.util import Pt
            from pptx.oxml.ns import qn
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import lxml.etree as etree

            # Positions extraites du PDF selon le type
            if slide_type == 'envoyer':  # Prez développée (21 slides), slide 4
                CARD_POSITIONS = {
                    'silos':           (0.059, 0.303, 0.212, 0.198),
                    'flotte':          (0.282, 0.303, 0.212, 0.198),
                    'tco':             (0.505, 0.303, 0.212, 0.198),
                    'sinistralite':    (0.729, 0.303, 0.212, 0.198),
                    'contrats':        (0.059, 0.521, 0.212, 0.198),
                    'fournisseurs':    (0.282, 0.521, 0.212, 0.198),
                    'electrification': (0.505, 0.521, 0.212, 0.198),
                    'conformite':      (0.729, 0.521, 0.212, 0.198),
                }
            else:  # Prez courte (11 slides), slide 3
                CARD_POSITIONS = {
                    'silos':           (0.059, 0.374, 0.212, 0.158),
                    'flotte':          (0.283, 0.374, 0.211, 0.158),
                    'tco':             (0.506, 0.374, 0.212, 0.158),
                    'sinistralite':    (0.730, 0.374, 0.211, 0.158),
                    'contrats':        (0.059, 0.553, 0.212, 0.159),
                    'fournisseurs':    (0.283, 0.553, 0.211, 0.159),
                    'electrification': (0.506, 0.553, 0.212, 0.159),
                    'conformite':      (0.730, 0.553, 0.211, 0.159),
                }

            W = prs.slide_width
            H = prs.slide_height

            for pp in pain_points:
                pos = CARD_POSITIONS.get(pp)
                if not pos:
                    continue
                l, t, w, h = pos
                shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    Emu(int(W * l)), Emu(int(H * t)),
                    Emu(int(W * w)), Emu(int(H * h)),
                )
                shape.fill.background()  # transparent
                shape.line.color.rgb = FLEETI_GREEN
                shape.line.width = Pt(2.5)

        # ── Personnalisations sur la slide de couverture (index 0) ────────────
        if idx == 0 and client_name:

            # Nom du client — 36pt foncé gras, aligné avec le texte du slide (5%)
            add_text_box(slide,
                left_pct=0.05, top_pct=0.725,
                width_pct=0.60, height_pct=0.10,
                text=client_name,
                font_size=36, bold=True, color=FLEETI_DARK,
            )

            # Infos flotte — 13pt vert, même alignement
            fleet_parts = []
            if vehicles and vehicles > 0:
                fleet_parts.append(f'{vehicles} véhicules')
            if vehicle_types:
                fleet_parts.append(' · '.join(VTYPE_SHORT[v] for v in vehicle_types))
            if fleet_parts:
                add_text_box(slide,
                    left_pct=0.05, top_pct=0.845,
                    width_pct=0.60, height_pct=0.07,
                    text='  ·  '.join(fleet_parts),
                    font_size=13, bold=False, color=FLEETI_GREEN,
                )

            # Pas de pain points sur la couverture — mis en avant sur slide 3

            # Logo client — slide 1 (tout en bas à droite)
            if logo_path:
                add_logo(slide, logo_path, left_pct=0.80, top_pct=0.88, width_pct=0.14)

    prs.save(output_path)
    print(f'Saved: {output_path} ({len(slide_files)} slides)')


def main():
    p = argparse.ArgumentParser()
    p.add_argument('--type', default='envoyer', choices=['envoyer', 'rdv'])
    p.add_argument('--client', default='')
    p.add_argument('--vehicles', type=int, default=0)
    p.add_argument('--vehicle-types', default='')
    p.add_argument('--pain-points', default='')
    p.add_argument('--logo-path', default='')
    p.add_argument('--sales', default='')
    p.add_argument('--output', required=True)
    args = p.parse_args()

    build_pptx(
        slide_type=args.type,
        client_name=args.client,
        vehicles=args.vehicles,
        vehicle_types=[v.strip() for v in args.vehicle_types.split(',') if v.strip()],
        pain_points=[v.strip() for v in args.pain_points.split(',') if v.strip()],
        logo_path=args.logo_path,
        sales_name=args.sales,
        output_path=args.output,
    )

if __name__ == '__main__':
    main()
